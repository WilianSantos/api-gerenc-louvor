from django_filters.rest_framework import DjangoFilterBackend
from drf_yasg import openapi
from drf_yasg.utils import swagger_auto_schema
from rest_framework import filters, status, viewsets
from rest_framework.response import Response
from rest_framework.views import APIView
from rest_framework.decorators import action
from datetime import timedelta
from dateutil.relativedelta import relativedelta
from django.db.models import Count
from django.utils import timezone
from pptx import Presentation
from pptx.util import Inches, Pt
from django.http import HttpResponse

from apps.accounts.models import Member
from apps.music.models import Music
from apps.music.utils import extract_lyrics_without_chords, is_chord
from apps.playlist.models import Playlist

from .models import LineupMember, PraiseLineup
from .serializers import LineupMemberSerializers, PraiseLineupSerializers


class PraiseLineupViewSet(viewsets.ModelViewSet):
    queryset = PraiseLineup.objects.all()
    serializer_class = PraiseLineupSerializers

    filter_backends = [
        DjangoFilterBackend,
        filters.OrderingFilter,
        filters.SearchFilter,
    ]
    ordering_fields = ["lineup_date", "lineup_event"]
    search_fields = ["lineup_date", "lineup_event"]

    @swagger_auto_schema(
        method='get',
        operation_description="Lista de escalas com links da playlist",
        responses={
            200: openapi.Response(
                description="Lista de escalas",
                schema=openapi.Schema(
                    type=openapi.TYPE_OBJECT,
                    properties={
                        "scales": openapi.Schema(
                            type=openapi.TYPE_ARRAY,
                            items=openapi.Schema(
                                type=openapi.TYPE_OBJECT,
                                properties={
                                    "id": openapi.Schema(type=openapi.TYPE_INTEGER),
                                    "lineup_event": openapi.Schema(type=openapi.TYPE_STRING),
                                    "lineup_date": openapi.Schema(type=openapi.TYPE_STRING, format="date"),
                                    "playlist": openapi.Schema(type=openapi.TYPE_INTEGER),
                                    "playlist_display": openapi.Schema(type=openapi.TYPE_STRING),
                                    "playlist_link_display": openapi.Schema(
                                        type=openapi.TYPE_ARRAY,
                                        items=openapi.Schema(type=openapi.TYPE_STRING)
                                    )
                                }
                            )
                        )
                    }
                )
            )
        }
    )
    @action(
        detail=False,
        methods=["get"],
        url_path="scales"
        
    )
    def get_scales(self, request):
        queryset = self.filter_queryset(self.get_queryset())  
        serializer = self.get_serializer(queryset, many=True)
        return Response({"scales": serializer.data}, status=status.HTTP_200_OK)

    @swagger_auto_schema(
        method='get',
        operation_description="Retorna total de escalas",
        responses={
            200: openapi.Response(
                description="Lista de escalas",
                schema=openapi.Schema(
                    type=openapi.TYPE_OBJECT,
                    properties={
                        "total": openapi.Schema(
                            type=openapi.TYPE_INTEGER,
                        )
                    }
                )
            )
        }
    )
    @action(
        detail=False,
        methods=["get"],
        url_path="total-scales"
        
    )
    def get_total_scales(self, request):
        scales = PraiseLineup.objects.all()
        
        return Response({"total": scales.count()}, status=status.HTTP_200_OK)
    
    @swagger_auto_schema(
    method='get',
    operation_description="Retorna as próximas 5 escalas dentro de 6 meses a partir da data atual.",
    responses={
            200: openapi.Response(
                description="Lista das próximas escalas",
                schema=openapi.Schema(
                    type=openapi.TYPE_OBJECT,
                    properties={
                        "next-scales": openapi.Schema(
                            type=openapi.TYPE_ARRAY,
                            items=openapi.Schema(
                                type=openapi.TYPE_OBJECT,
                                properties={
                                    "event": openapi.Schema(type=openapi.TYPE_STRING),
                                    "date": openapi.Schema(type=openapi.TYPE_STRING, format="date")
                                }
                            )
                        )
                    }
                )
            )
        }
    )
    @action(
        detail=False,
        methods=["get"],
        url_path="next-scales"
    )
    @action(
        detail=False,
        methods=["get"],
        url_path="next-scales"
        
    )
    def get_next_scales(self, request):
        date_now = timezone.now()
        date_next_months = date_now + relativedelta(months=6)
        

        lineups = PraiseLineup.objects.filter(
            lineup_date__range=(date_now, date_next_months)
        )[:5]
        
        return Response({"next-scales": [
           {
            "event": scale.lineup_event,
            "date": scale.lineup_date
           }
            for scale in lineups
        ]}, status=status.HTTP_200_OK)
    
    @swagger_auto_schema(
    method='get',
    operation_description="Retorna as últimas 5 escalas dos 6 meses anteriores à data atual.",
    responses={
            200: openapi.Response(
                description="Lista das escalas anteriores",
                schema=openapi.Schema(
                    type=openapi.TYPE_OBJECT,
                    properties={
                        "previous-scales": openapi.Schema(
                            type=openapi.TYPE_ARRAY,
                            items=openapi.Schema(
                                type=openapi.TYPE_OBJECT,
                                properties={
                                    "event": openapi.Schema(type=openapi.TYPE_STRING),
                                    "date": openapi.Schema(type=openapi.TYPE_STRING, format="date")
                                }
                            )
                        )
                    }
                )
            )
        }
    )
    @action(
        detail=False,
        methods=["get"],
        url_path="previous-scales"
        
    )
    def get_previous_scales(self, request):
        date_now = timezone.now()
        date_previous_months = date_now - relativedelta(months=6)
        

        lineups = PraiseLineup.objects.filter(
            lineup_date__range=(date_previous_months, date_now)
        )[:5]
        
        return Response({"previous-scales": [
           {
            "event": scale.lineup_event,
            "date": scale.lineup_date
           }
            for scale in lineups
        ]}, status=status.HTTP_200_OK)


class LineupMemberViewSet(viewsets.ModelViewSet):
    queryset = LineupMember.objects.all()
    serializer_class = LineupMemberSerializers
    pagination_class = None

    filter_backends = [
        DjangoFilterBackend,
        filters.OrderingFilter,
        filters.SearchFilter,
    ]
    ordering_fields = ["lineup__lineup_date", "lineup__lineup_event"]
    search_fields = ["member__name", "function__function_name"]


class ScaleHistoryViewSet(APIView):
    @swagger_auto_schema(
        operation_description="Rota para retornar o histórico de escalas do membro logado.",
        responses={
            200: openapi.Response(
                description="Lista de escalas",
                schema=openapi.Schema(
                    type=openapi.TYPE_OBJECT,
                    properties={
                        "scales": openapi.Schema(
                            type=openapi.TYPE_ARRAY,
                            items=openapi.Schema(
                                type=openapi.TYPE_OBJECT,
                                properties={
                                    "id": openapi.Schema(type=openapi.TYPE_INTEGER),
                                    "lineup_date": openapi.Schema(
                                        type=openapi.TYPE_STRING,
                                        format=openapi.FORMAT_DATE,
                                    ),
                                    "lineup_event": openapi.Schema(
                                        type=openapi.TYPE_STRING
                                    ),
                                    "function_name": openapi.Schema(
                                        type=openapi.TYPE_STRING
                                    ),
                                },
                            ),
                        )
                    },
                ),
            )
        },
    )
    def get(self, request):
        member_id = Member.objects.get(user=request.user).id
        scales = LineupMember.objects.filter(member_id=member_id)
        scales = scales.values(
            "id",
            "lineup__lineup_date",
            "lineup__lineup_event",
            "function__function_name",
        ).order_by("-lineup__lineup_date")[:10]
        return Response(
            {
                "scales": [
                    {
                        "id": scale["id"],
                        "lineup_date": scale["lineup__lineup_date"],
                        "lineup_event": scale["lineup__lineup_event"],
                        "function_name": scale["function__function_name"],
                    }
                    for scale in scales
                ]
            },
            status=status.HTTP_200_OK,
        )


class SlideGeneratorView(APIView):
    @swagger_auto_schema(
        operation_description="Gera e retorna um arquivo .pptx com as músicas da playlist (sem acordes).",
        request_body=openapi.Schema(
            type=openapi.TYPE_OBJECT,
            required=["playlist_id"],
            properties={
                "playlist_id": openapi.Schema(
                    type=openapi.TYPE_INTEGER,
                    description="ID da playlist cujas músicas serão usadas para gerar os slides"
                )
            }
        ),
        responses={
            200: openapi.Response(
                description="Arquivo .pptx com os slides das músicas",
                schema=openapi.Schema(
                    type=openapi.TYPE_STRING,
                    format='binary'  # <- indica que é um arquivo
                )
            )
        }
    )
    def post(self, request):
        playlist_id = request.data.get("playlist_id")
        if not playlist_id:
            return Response({"playlist_id": "ID não foi passado"}, status=400)
        
        try:
            playlist = Playlist.objects.get(id=playlist_id)
        except Playlist.DoesNotExist:
            return Response({"detail": "Playlist não encontrada"}, status=404)

        through_model = playlist.music.through  # acesso à tabela M2M

        musics = (
            through_model.objects
            .values("music")
        )
        music_ids = [music["music"] for music in musics]
        music_list = Music.objects.filter(id__in=music_ids)
        # Cria uma nova apresentação
        prs = Presentation()
        title_slide_layout = prs.slide_layouts[5]  # Layout em branco

        # Gera slides para cada música
        for music in music_list:
            # Extrai letras sem acordes
            music_text = extract_lyrics_without_chords(music.music_text)
            lines = music_text.strip().split('\n')  # cada linha separada

            # Controla o número de slides para esta música
            slide_count = 0
            lines_on_current_slide = 0

            # Cria o primeiro slide para a música
            slide = prs.slides.add_slide(title_slide_layout)
            slide_count += 1
            
            # Adiciona título no primeiro slide
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
            title_tf = title_box.text_frame
            title_tf.text = f"{music.music_title} - {music.author}"
            title_tf.paragraphs[0].font.size = Pt(32)
            title_tf.paragraphs[0].font.bold = True

            # Adiciona texto do slide
            textbox = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(5.5))
            tf = textbox.text_frame
            tf.word_wrap = True

            # Itera sobre todas as linhas da música
            for line in lines:
                # Se já foram adicionadas 5 linhas, cria um novo slide
                if lines_on_current_slide >= 5:
                    slide = prs.slides.add_slide(title_slide_layout)
                    slide_count += 1
                    
                    # Reseta o textbox para o novo slide
                    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(5.5))
                    tf = textbox.text_frame
                    tf.word_wrap = True
                    
                    # Reseta o contador de linhas
                    lines_on_current_slide = 0

                # Adiciona a linha ao slide atual
                if line.strip():
                    text = ''
                    for t in line.split():
                        if not is_chord(t):
                            text += t + ' '
                    p = tf.add_paragraph()
                    p.text = text.strip()
                    p.font.size = Pt(30)
                    lines_on_current_slide += 1

        # Salva a apresentação
        response = HttpResponse(content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
        response['Content-Disposition'] = 'attachment; filename="culto_slides.pptx"'
        prs.save(response)
        return response
    
